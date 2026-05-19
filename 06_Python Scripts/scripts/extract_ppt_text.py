"""Extract text and speaker notes from all PPTX files in pptx/ → input_md/"""
from pptx import Presentation
from pathlib import Path

PPTX_DIR = Path("pptx")
OUTPUT_DIR = Path("input_md")
OUTPUT_DIR.mkdir(exist_ok=True)

def safe_filename(name):
    invalid_chars = '<>:"/\\|?*'
    for ch in invalid_chars:
        name = name.replace(ch, "_")
    return name

def extract_text_from_shape(shape):
    texts = []
    if hasattr(shape, "text") and shape.text and shape.text.strip():
        texts.append(shape.text.strip())
    if shape.has_table:
        for row in shape.table.rows:
            row_texts = []
            for cell in row.cells:
                if cell.text.strip():
                    row_texts.append(cell.text.strip())
            if row_texts:
                texts.append(" | ".join(row_texts))
    if hasattr(shape, "shapes"):
        for sub_shape in shape.shapes:
            texts.extend(extract_text_from_shape(sub_shape))
    return texts

def extract_pptx(pptx_path):
    prs = Presentation(pptx_path)
    md_lines = []

    md_lines.append(f"# {pptx_path.stem}")
    md_lines.append("")
    md_lines.append(f"- 总页数：{len(prs.slides)}")
    md_lines.append("")

    for idx, slide in enumerate(prs.slides, start=1):
        md_lines.append(f"## 第 {idx} 页")
        md_lines.append("")

        # --- Slide text ---
        md_lines.append("### 原始页面文字")
        md_lines.append("")

        slide_texts = []
        for shape in slide.shapes:
            slide_texts.extend(extract_text_from_shape(shape))

        if slide_texts:
            md_lines.append("\n\n".join(slide_texts))
        else:
            md_lines.append("本页未提取到文字，可能是图片型页面，请补充截图 OCR 或视觉描述。")

        md_lines.append("")

        # --- Speaker notes ---
        md_lines.append("### 原始备注 (Speaker Notes)")
        md_lines.append("")
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                md_lines.append(notes)
            else:
                md_lines.append("（无备注）")
        else:
            md_lines.append("（无备注）")

        md_lines.append("")
        md_lines.append("---")
        md_lines.append("")

    output_path = OUTPUT_DIR / (safe_filename(pptx_path.stem) + ".md")
    output_path.write_text("\n".join(md_lines), encoding="utf-8")
    print(f"已生成：{output_path}")

def main():
    pptx_files = sorted(PPTX_DIR.glob("*.pptx"))
    if not pptx_files:
        print("pptx/ 文件夹中没有找到 .pptx 文件。")
        return

    for pptx_file in pptx_files:
        print(f"处理: {pptx_file.name}")
        extract_pptx(pptx_file)

    print(f"\n完成，共处理 {len(pptx_files)} 个文件。")

if __name__ == "__main__":
    main()
