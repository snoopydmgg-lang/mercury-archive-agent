"""Read speech MD files and write speaker notes back into PPTX files."""
from pptx import Presentation
from pathlib import Path
import re

PPTX_DIR = Path("pptx")
OUTPUT_DIR = Path("output")
NOTES_DIR = Path("output_pptx")
NOTES_DIR.mkdir(exist_ok=True)

md_files = {
    "1": "1_全面机制升维_逐页演讲稿.md",
    "2": "2_机制突破_逐页演讲稿.md",
    "3": "3_疗效升维_逐页演讲稿.md",
    "4": "4_安全革新_逐页演讲稿.md",
}

def parse_speeches(md_path):
    """Parse MD file and return dict of page_num -> speech text."""
    text = md_path.read_text(encoding="utf-8")

    # Split into pages — handle both "## 第1页" and "## 第 1 页" and "## 第1页：xxx"
    page_blocks = re.split(r'\n## 第\s*(\d+)\s*页', text)
    speeches = {}

    # page_blocks[0]=preamble, [1]=page_num, [2]=content, [3]=page_num, ...
    for i in range(1, len(page_blocks), 2):
        page_num = int(page_blocks[i])
        content = page_blocks[i + 1] if i + 1 < len(page_blocks) else ""

        # Find ### 建议演讲稿 or ### 2. 建议演讲稿 section
        # Match from this header to the next ### header
        m_start = re.search(r'^###\s+(?:2\.\s*)?建议演讲稿\s*$', content, re.MULTILINE)
        if not m_start:
            continue

        start_pos = m_start.end()
        remaining = content[start_pos:]

        # Find the next ### section header as end marker
        m_end = re.search(r'^###\s+', remaining, re.MULTILINE)
        if m_end:
            speech = remaining[:m_end.start()].strip()
        else:
            speech = remaining.strip()

        # Also stop at --- (page separator)
        speech = re.split(r'\n---\n', speech)[0].strip()

        if speech:
            speeches[page_num] = speech

    return speeches

def main():
    for prefix, md_name in md_files.items():
        md_path = OUTPUT_DIR / md_name
        if not md_path.exists():
            print(f"SKIP: {md_name} not found")
            continue

        speeches = parse_speeches(md_path)
        print(f"{md_name}: parsed {len(speeches)} speeches")

        pptx_files = [f for f in PPTX_DIR.glob("*.pptx") if f.stem.startswith(prefix)]
        if not pptx_files:
            print(f"  No PPTX found starting with {prefix}")
            continue

        pptx_path = pptx_files[0]
        prs = Presentation(str(pptx_path))
        print(f"  Source: {pptx_path.name}")

        updated = 0
        for idx, slide in enumerate(prs.slides, start=1):
            speech = speeches.get(idx)
            if not speech:
                continue
            notes_slide = slide.notes_slide
            tf = notes_slide.notes_text_frame
            tf.clear()
            tf.text = speech
            updated += 1

        out_path = NOTES_DIR / pptx_path.name
        prs.save(str(out_path))
        print(f"  Saved: {out_path.name} ({updated}/{len(prs.slides)} slides updated)")

    print(f"\nDone. Output in {NOTES_DIR}/")

if __name__ == "__main__":
    main()
